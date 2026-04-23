from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Premium color palette
PRIMARY_GREEN = RGBColor(4, 232, 141)      # Neon green
PRIMARY_DARK = RGBColor(2, 192, 119)       # Dark green
ACCENT_BLUE = RGBColor(100, 150, 255)     # Modern blue
ACCENT_PURPLE = RGBColor(150, 100, 255)   # Modern purple
BG_DARK = RGBColor(15, 15, 35)            # Very dark navy
BG_CARD = RGBColor(25, 30, 50)            # Card background
TEXT_WHITE = RGBColor(255, 255, 255)      # White
TEXT_LIGHT = RGBColor(220, 220, 240)      # Light grey
TEXT_MUTED = RGBColor(150, 160, 190)      # Muted
ACCENT_ORANGE = RGBColor(255, 150, 80)    # Warm accent

def create_gradient_bg(slide, color1, color2):
    """Create gradient-like effect with overlapping shapes"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK

def add_decorative_element(slide, x, y, width, height, color, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    """Add decorative shape"""
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    return shape

def add_title_slide(prs):
    """Premium title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK

    # Decorative shapes (top right)
    shape1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(-0.5), Inches(4), Inches(4))
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = RGBColor(4, 232, 141)
    shape1.line.color.rgb = PRIMARY_GREEN
    shape1.line.width = Pt(2)
    shape1.rotation = 45

    # Decorative shapes (bottom left)
    shape2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(-1), Inches(5.5), Inches(3.5), Inches(3.5))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = ACCENT_BLUE
    shape2.line.color.rgb = ACCENT_BLUE
    shape2.line.width = Pt(2)
    shape2.rotation = 30

    # Header bar with gradient effect
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_GREEN
    header.line.color.rgb = PRIMARY_GREEN

    # University name (white on green)
    univ_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.8))
    univ_frame = univ_box.text_frame
    p = univ_frame.paragraphs[0]
    p.text = "GANPAT UNIVERSITY"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BG_DARK
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Calibri"

    # Main title box
    title_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(2.2), Inches(7), Inches(1.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = BG_CARD
    title_shape.line.color.rgb = PRIMARY_GREEN
    title_shape.line.width = Pt(3)

    # Project title
    title_box = slide.shapes.add_textbox(Inches(1.7), Inches(2.4), Inches(6.6), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = "CodeHub"
    p.font.size = Pt(80)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_GREEN
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Calibri"

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(0.7))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = "Duolingo-Style Coding Learning Platform"
    p.font.size = Pt(28)
    p.font.color.rgb = TEXT_LIGHT
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Calibri"

    # Guide info box
    guide_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(5.2), Inches(8), Inches(1))
    guide_shape.fill.solid()
    guide_shape.fill.fore_color.rgb = RGBColor(35, 40, 70)
    guide_shape.line.color.rgb = ACCENT_BLUE
    guide_shape.line.width = Pt(2)

    guide_box = slide.shapes.add_textbox(Inches(1.2), Inches(5.3), Inches(7.6), Inches(0.8))
    guide_frame = guide_box.text_frame
    p = guide_frame.paragraphs[0]
    p.text = "Internal Guide: Prof. Chirag Patel"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Calibri"

    # Student info
    student_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.8))
    student_frame = student_box.text_frame
    p = student_frame.paragraphs[0]
    p.text = "B.Tech Final Semester Project  |  Abhay Prajapati"
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_MUTED
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Calibri"

def add_premium_content_slide(prs, title, emoji, content_items):
    """Premium content slide with better design"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK

    # Accent bar on left
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = PRIMARY_GREEN
    accent_bar.line.color.rgb = PRIMARY_GREEN

    # Title background
    title_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.3), Inches(9), Inches(0.9))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = BG_CARD
    title_bg.line.color.rgb = PRIMARY_GREEN
    title_bg.line.width = Pt(2)

    # Title with emoji
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(8.6), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = f"{emoji}  {title}"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_GREEN
    p.font.name = "Calibri"

    # Content box
    content_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(5.7))
    content_bg.fill.solid()
    content_bg.fill.fore_color.rgb = BG_CARD
    content_bg.line.color.rgb = RGBColor(50, 60, 100)
    content_bg.line.width = Pt(1)

    # Content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(8), Inches(5.3))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, item in enumerate(content_items):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(22)
        p.font.color.rgb = TEXT_LIGHT
        p.font.name = "Calibri"
        p.space_before = Pt(8)
        p.space_after = Pt(8)

def add_two_col_premium(prs, title, emoji, left_items, right_items):
    """Two column premium slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK

    # Accent bar
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT_BLUE
    accent_bar.line.color.rgb = ACCENT_BLUE

    # Title background
    title_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.3), Inches(9), Inches(0.9))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = BG_CARD
    title_bg.line.color.rgb = ACCENT_BLUE
    title_bg.line.width = Pt(2)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(8.6), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = f"{emoji}  {title}"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.font.name = "Calibri"

    # Left column box
    left_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(4.4), Inches(5.7))
    left_bg.fill.solid()
    left_bg.fill.fore_color.rgb = RGBColor(30, 35, 60)
    left_bg.line.color.rgb = PRIMARY_GREEN
    left_bg.line.width = Pt(2)

    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(3.8), Inches(5.3))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    for i, item in enumerate(left_items):
        if i > 0:
            left_frame.add_paragraph()
        p = left_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_LIGHT
        p.font.name = "Calibri"
        p.space_before = Pt(6)
        p.space_after = Pt(6)

    # Right column box
    right_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.1), Inches(1.4), Inches(4.4), Inches(5.7))
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = RGBColor(30, 35, 60)
    right_bg.line.color.rgb = ACCENT_PURPLE
    right_bg.line.width = Pt(2)

    right_box = slide.shapes.add_textbox(Inches(5.4), Inches(1.6), Inches(3.8), Inches(5.3))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    for i, item in enumerate(right_items):
        if i > 0:
            right_frame.add_paragraph()
        p = right_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_LIGHT
        p.font.name = "Calibri"
        p.space_before = Pt(6)
        p.space_after = Pt(6)

# Slide 1: Title
add_title_slide(prs)

# Slide 2: Overview
add_premium_content_slide(prs, "Project Overview", "🎯", [
    "Interactive Duolingo-style coding learning platform",
    "Bite-sized lessons with multiple-choice quizzes",
    "Gamification: XP rewards, streaks, hearts system",
    "Real-time leaderboard with competitive rankings",
    "User authentication & progress tracking",
    "Neon green UI inspired by Duolingo"
])

# Slide 3: Tech Stack
add_two_col_premium(prs, "Tech Stack", "⚙️", [
    "FRONTEND:",
    "  Angular 17+",
    "  TypeScript",
    "  Standalone Components",
    "  Modern CSS (Dark Theme)",
    "",
    "DATABASE:",
    "  MySQL (XAMPP)",
    "  phpMyAdmin",
    "  5 core tables"
], [
    "BACKEND:",
    "  PHP 7.4+",
    "  REST API",
    "  CORS Enabled",
    "  10 Endpoints",
    "",
    "HTTP CLIENT:",
    "  Angular HttpClient",
    "  localStorage Session",
    "  JSON API"
])

# Slide 4: Features
add_premium_content_slide(prs, "Key Features", "✨", [
    "User registration & authentication with password hashing",
    "Course library with structured lessons & content",
    "Interactive quiz system with instant feedback",
    "Hearts system (5 max, lose 1 per wrong answer)",
    "XP rewards & consecutive day streak tracking",
    "Live leaderboard showing top 10 users",
    "User profiles with statistics & badges"
])

# Slide 5: Database
add_premium_content_slide(prs, "Database Schema", "🗄️", [
    "users: id, name, email, password, xp, streak, hearts, last_active",
    "",
    "courses: id, title, description, icon, color",
    "",
    "lessons: id, course_id, title, content, xp_reward, order_num",
    "",
    "quiz_questions: id, lesson_id, question, options (a-d), correct_option",
    "",
    "user_progress: id, user_id, lesson_id, completed, score, completed_at"
])

# Slide 6: Pages
add_premium_content_slide(prs, "8 Application Pages", "📱", [
    "Home - Public landing page",
    "Login / Register - User authentication forms",
    "Dashboard - Winding lesson path with progress tracking",
    "Courses - Grid view of all available courses",
    "Lessons - Lesson list with content preview per course",
    "Quiz - Interactive questions with instant feedback",
    "Leaderboard - Top 10 rankings by XP points",
    "Profile - User stats, level, badges, achievements"
])

# Slide 7: Design
add_premium_content_slide(prs, "Design System", "🎨", [
    "Color Palette:",
    "  Primary: Neon Green #04E88D  |  Background: Dark #0F0F23",
    "  Cards: #191E32 with depth effect",
    "  Danger: #FF4B4B (hearts/errors)  |  Accent: #6496FF (blue)",
    "",
    "Typography: Calibri font with modern hierarchy",
    "3D buttons with shadow effect on interaction",
    "Rounded rectangles for modern aesthetic"
])

# Slide 8: Hearts
add_premium_content_slide(prs, "Hearts System", "❤️", [
    "5 hearts maximum per user at any time",
    "Lose 1 heart per incorrect quiz answer",
    "Quiz is blocked when user has 0 hearts",
    "Daily auto-refill of hearts after 24 hours",
    "Hearts synced with backend database",
    "Encourages daily engagement without frustration",
    "Matches Duolingo's proven engagement model"
])

# Slide 9: API
add_two_col_premium(prs, "Backend API Endpoints", "🔌", [
    "AUTHENTICATION:",
    "  POST /register.php",
    "  POST /login.php",
    "",
    "CONTENT:",
    "  GET /get_courses.php",
    "  GET /get_lessons.php",
    "  GET /get_quiz.php",
    "",
    "PROGRESS:",
    "  POST /save_progress.php",
    "  GET /get_profile.php"
], [
    "HEARTS:",
    "  POST /update_hearts.php",
    "",
    "SOCIAL:",
    "  GET /get_leaderboard.php",
    "",
    "BASE URL:",
    "  http://localhost/",
    "  codehub-api/",
    "",
    "All endpoints return JSON",
    "CORS enabled for frontend"
])

# Slide 10: Journey
add_premium_content_slide(prs, "User Journey", "🚀", [
    "Register/Login → Session stored in localStorage with codehub_user key",
    "View Dashboard → Winding path shows progress (done/active/locked nodes)",
    "Select Course → Browse courses in grid view with descriptions",
    "Choose Lesson → View content & XP reward information",
    "Take Quiz → Answer multiple-choice questions with feedback",
    "Earn XP & Streak → Profile updates in real-time with statistics",
    "Climb Leaderboard → Compete globally with other users"
])

# Slide 11: Status
add_premium_content_slide(prs, "Implementation Status", "✓", [
    "✓ All 8 pages built with Angular 17+ standalone components",
    "✓ All PHP REST endpoints functional & tested",
    "✓ MySQL database with sample data loaded",
    "✓ Dark theme + neon green branding fully applied",
    "✓ Login/Register with secure session storage",
    "✓ Quiz system with hearts working perfectly",
    "✓ Leaderboard & profile pages complete"
])

# Slide 12: How to Run
add_premium_content_slide(prs, "How to Run the Project", "🖥️", [
    "PREREQUISITES: Node.js 18+, XAMPP (Apache + MySQL), VS Code",
    "",
    "START BACKEND:",
    "  1. Open XAMPP Control Panel",
    "  2. Start Apache & MySQL services",
    "  3. Verify: http://localhost/phpmyadmin",
    "",
    "START FRONTEND:",
    "  Run: ng serve (port 4200)",
    "  Open: http://localhost:4200"
])

# Slide 13: Challenges
add_premium_content_slide(prs, "Challenges & Solutions", "🔧", [
    "Real-time Progress → Used POST API calls after quiz completion",
    "Hearts Sync → Store in localStorage + backend database",
    "Responsive Dark UI → CSS custom properties & media queries",
    "CORS with PHP → Proper headers in config.php",
    "User Sessions → localStorage key: codehub_user",
    "Leaderboard Ranking → Implemented XP-based sorting",
    "Mobile Responsiveness → Tested on all device sizes"
])

# Slide 14: Future
add_premium_content_slide(prs, "Future Enhancements", "🌟", [
    "Mobile app version (React Native / Flutter)",
    "Real-time hints & code explanations in lessons",
    "Social features (follow friends, chat, discussion)",
    "Certificate generation on course completion",
    "Admin dashboard for course/lesson management",
    "Dark/Light theme toggle",
    "Advanced achievement badges system"
])

# Slide 15: Closing
closing = prs.slides.add_slide(prs.slide_layouts[6])
background = closing.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = BG_DARK

# Decorative shapes
shape1 = closing.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(-0.5), Inches(4), Inches(4))
shape1.fill.solid()
shape1.fill.fore_color.rgb = PRIMARY_GREEN
shape1.line.color.rgb = PRIMARY_GREEN
shape1.rotation = 45

shape2 = closing.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(-1), Inches(5), Inches(3), Inches(3))
shape2.fill.solid()
shape2.fill.fore_color.rgb = ACCENT_BLUE
shape2.line.color.rgb = ACCENT_BLUE
shape2.rotation = 30

# Main box
main_bg = closing.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.5), Inches(7), Inches(4.5))
main_bg.fill.solid()
main_bg.fill.fore_color.rgb = BG_CARD
main_bg.line.color.rgb = PRIMARY_GREEN
main_bg.line.width = Pt(3)

# Thank you
thank_box = closing.shapes.add_textbox(Inches(1.7), Inches(1.8), Inches(6.6), Inches(1))
thank_frame = thank_box.text_frame
p = thank_frame.paragraphs[0]
p.text = "Thank You!"
p.font.size = Pt(72)
p.font.bold = True
p.font.color.rgb = PRIMARY_GREEN
p.alignment = PP_ALIGN.CENTER
p.font.name = "Calibri"

# Questions
q_box = closing.shapes.add_textbox(Inches(1.7), Inches(3), Inches(6.6), Inches(1))
q_frame = q_box.text_frame
p = q_frame.paragraphs[0]
p.text = "Questions & Discussion"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = TEXT_LIGHT
p.alignment = PP_ALIGN.CENTER
p.font.name = "Calibri"

# Tagline
tag_box = closing.shapes.add_textbox(Inches(1.7), Inches(4.2), Inches(6.6), Inches(0.8))
tag_frame = tag_box.text_frame
p = tag_frame.paragraphs[0]
p.text = "Learn. Play. Master Code."
p.font.size = Pt(26)
p.font.color.rgb = ACCENT_BLUE
p.alignment = PP_ALIGN.CENTER
p.font.italic = True
p.font.name = "Calibri"

# Save
output_path = r"e:\College 8\codehub\CodeHub_Presentation.pptx"
prs.save(output_path)
print(f"Presentation created: {output_path}")
